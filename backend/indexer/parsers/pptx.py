"""PPTX raw-content parser."""

from __future__ import annotations

from pptx import Presentation

from backend.indexer.parsers import ContentSegment, DiscoveredFile, ParsedContent


def parse_pptx(file: DiscoveredFile) -> ParsedContent:
    """Extract raw slide text from a PPTX."""

    presentation = Presentation(file.path)
    slides: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                slide_text.append(text)
        slides.append("\n".join(slide_text))

    segments = [
        ContentSegment(locator=f"s{index}", text=slide_text)
        for index, slide_text in enumerate(slides, start=1)
    ]
    return ParsedContent(text="\n\n".join(slides), segments=segments)
